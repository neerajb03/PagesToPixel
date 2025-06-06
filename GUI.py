import os
import streamlit as st
import requests
import requests.exceptions
import time
import google.generativeai as genai
from dotenv import load_dotenv
import fitz  # PyMuPDF for PDF text extraction
from docx import Document  # For DOCX file handling
from pptx import Presentation  # For PPTX file handling

def show():
    # Add animation effect
    st.markdown('<div class="fade-in">', unsafe_allow_html=True)
    
    # Initialize variables
    num_lines = None
    num_paragraphs = None
    custom_prompt = None

    # Load environment variables from .env file
    load_dotenv()

    # Configure API keys
    GENAI_API_KEY = st.secrets["GENAI_API_KEY"]
    HEYGEN_API_KEY = st.secrets["HEYGEN_API_KEY"]

    genai.configure(api_key=GENAI_API_KEY)

    # Minimal header
    st.title('Document to Video')
    st.markdown("##### Create AI-narrated videos from your documents")

    # Function to extract text from different file formats
    def extract_text(file_path, file_extension):
        try:
            if file_extension == "pdf":
                with fitz.open(file_path) as pdf_file:
                    return "\n".join([page.get_text() for page in pdf_file])
            elif file_extension == "docx":
                doc = Document(file_path)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif file_extension == "pptx":
                presentation = Presentation(file_path)
                return "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
            return None
        except Exception as e:
            st.error(f"Error: {str(e)}")
            return None

    # Function to generate text summary
    def generate_summary(file_text, summary_type, num_lines, num_paragraphs, custom_prompt):
        try:
            with st.spinner("Generating summary..."):
                model = genai.GenerativeModel("gemini-1.5-flash")
            
                if summary_type == "Number of Lines":
                    prompt = f"Summarize the following content in {num_lines} lines:"
                elif summary_type == "Number of Paragraphs":
                    prompt = f"Summarize the following content in {num_paragraphs} paragraphs:"
                else:
                    prompt = f"{custom_prompt}"
            
                response = model.generate_content([prompt, file_text])
                return response.text if hasattr(response, 'text') else None
        except Exception as e:
            st.error(f"Error: {str(e)}")
            return None

    # Function to generate a talking avatar video
    def generate_talking_avatar(audio_text):
        try:
            with st.spinner("Creating avatar video..."):
                url = "https://api.heygen.com/v2/video/generate"
            
                payload = {
                    "caption": False,
                    "title": "Generated Video",
                    "dimension": {"width": 1280, "height": 720},
                    "video_inputs": [
                        {
                            "character": {
                                "type": "avatar",
                                "avatar_id": "Timothy_sitting_office_front",
                                "scale": 1,
                                "avatar_style": "normal",
                                "offset": {"x": 0, "y": 0}
                            },
                            "voice": {
                                "type": "text",
                                "voice_id": "01d674cfd32b4728a3fddd21b7e7d543",
                                "input_text": audio_text
                            }
                        }
                    ]
                }
            
                headers = {
                    "accept": "application/json",
                    "content-type": "application/json",
                    "x-api-key": HEYGEN_API_KEY
                }
            
                response = requests.post(url, json=payload, headers=headers)
            
                if response.status_code == 200:
                    data = response.json().get("data", {})
                    return data.get("video_id")
            
                st.error(f"Error: {response.status_code}")
                return None
        except Exception as e:
            st.error(f"Error: {str(e)}")
            return None

    def get_video_url(video_id):
        try:
            with st.spinner("Processing video..."):
                url = f"https://api.heygen.com/v1/video_status.get?video_id={video_id}"
                headers = {"accept": "application/json", "x-api-key": HEYGEN_API_KEY}

                progress_bar = st.progress(0)
                max_attempts = 30  # Increased max attempts for longer processing
                wait_time = 20  # Time to wait between API calls

                last_status = None  # Store the last displayed status

                for attempt in range(max_attempts):
                    progress_value = min(0.95 * attempt / max_attempts, 0.95)
                    progress_bar.progress(progress_value)

                    try:
                        response = requests.get(url, headers=headers)
                        if response.status_code == 200:
                            data = response.json().get("data", {})
                            status = data.get("status", "")

                            # Display status only if it changes
                            if status != last_status:
                                if status == "processing":
                                    st.info(f"This may take a few minutes...")
                                elif status == "completed":
                                    progress_bar.progress(1.0)
                                    st.success("Video generation complete!")
                                    return data.get("video_url")
                                elif status == "failed":
                                    st.error("Video generation failed.")
                                    return None

                            last_status = status  # Update last status

                    except requests.exceptions.RequestException as e:
                        st.warning(f"API request issue: {str(e)}. Retrying...")

                    time.sleep(wait_time)

                st.warning("Maximum attempts reached. The video may still be processing.")
                st.info(f"You can check back later using this video ID: {video_id}")
                return None

        except Exception as e:
            st.error(f"Error in video retrieval: {str(e)}")
            return None


    # Create two columns for upload and results
    col1, col2 = st.columns([1, 1])
    
    #with col1:
        # Enhanced card for upload with new color scheme - REDUNDANT
        #st.markdown("""
        #<div class="card">
       #     <h4>Upload Document</h4>
        #""", unsafe_allow_html=True)
        
        # Upload file
    uploaded_file = st.file_uploader("Select a file", type=["pdf", "docx", "pptx"])
    
    if uploaded_file is not None:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        file_path = f"temp_file.{file_extension}"
    
        with open(file_path, "wb") as f:
            f.write(uploaded_file.read())
    
        st.success(f"✓ {uploaded_file.name}")
        
        # Summary options with improved styling and better contrast
        #with col2:
        st.markdown('<h5>Summary Options</h5>', unsafe_allow_html=True)
        summary_type = st.radio("Type:", ("Number of Lines", "Number of Paragraphs", "Custom Prompt"))

        if summary_type == "Number of Lines":   
            num_lines = st.text_area("Number of lines",placeholder="Enter number of lines here...",height=68)
        elif summary_type == "Number of Paragraphs":
            num_paragraphs = st.text_area("Number of paragraphs",placeholder="Enter number of paragraphs here...", height=68)
        else:
            custom_prompt = st.text_area("Custom prompt:", "Summarize this document concisely.", height=100)
    
    st.markdown("</div>", unsafe_allow_html=True)
    
    # Avatar options if file uploaded
    if uploaded_file is not None:
        #st.markdown(""" Unnessecary card
        #<div class="card">
        #    <h4>Generate</h4>
        #""", unsafe_allow_html=True)
        
        # Enhanced generate button with better contrast
        if st.button("Create Video", key="generate_btn", use_container_width=True):
            file_text = extract_text(file_path, file_extension)
        
            if file_text:
                # Use right column for results
                #with col2:
                st.markdown("""
                <div class="card">
                    <h4>Output</h4>
                """, unsafe_allow_html=True)
                
                audio_summary = generate_summary(file_text, summary_type, num_lines, num_paragraphs, custom_prompt)
            
                if audio_summary:
                    st.markdown('<h5>Summary:</h5>', unsafe_allow_html=True)
                    st.markdown(f'<div style="background-color:#f0f4f8; padding:10px; border-radius:8px; margin-bottom:15px; border: 1px solid #cbd5e0; color: #2d3748;">{audio_summary}</div>', unsafe_allow_html=True)
                
                    video_id = generate_talking_avatar(audio_summary)
                
                    if video_id:
                        video_url = get_video_url(video_id)
                    
                        if video_url:
                            st.success("✓ Video ready")
                            st.video(video_url)
                            st.markdown(f'<a href="{video_url}" target="_blank" style="display: inline-block; padding: 6px 12px; background-color: #2b6cb0; color: white; text-decoration: none; border-radius: 4px; font-weight: 500;">Download Video</a>', unsafe_allow_html=True)
                        else:
                            st.info("Video processing. Check back in a moment.")
                    else:
                        st.error("Couldn't generate video.")
                else:
                    st.error("Couldn't generate summary.")
                
                st.markdown("</div>", unsafe_allow_html=True)
            else:
                st.error("Couldn't extract text from file.")
        
        st.markdown("</div>", unsafe_allow_html=True)
    
    # Right column with enhanced instructions if no file uploaded
    #with col2:
    if uploaded_file is None:
        st.markdown("""
        <div class="card">
            <h4>How It Works</h4>
            <ol style="padding-left: 1.2rem; margin-bottom: 0;">
                <li>Upload document</li>
                <li>Choose summary options</li>
                <li>Generate video</li>
                <li>Download or share</li>
            </ol>
        </div>
        
        <div class="card">
            <h4>Supported Files</h4>
            <p>PDF, DOCX, and PPTX</p>
        </div>
        """, unsafe_allow_html=True)
    
    # End animation div
    st.markdown('</div>', unsafe_allow_html=True)